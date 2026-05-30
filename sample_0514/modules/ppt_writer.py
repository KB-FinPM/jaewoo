from typing import List

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from modules.schemas import ScreenPlanItem


def _set_text(shape, text: str, font_size: int = 14, bold: bool = False):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = "Malgun Gothic"


def _add_header(slide, item: ScreenPlanItem):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.55))
    _set_text(box, f"요구사항ID: {item.requirement_id}    화면번호: {item.screen_no}    화면명: {item.screen_name}", font_size=18, bold=True)


def _add_image_placeholder(slide, item: ScreenPlanItem):
    rect = slide.shapes.add_shape(1, Inches(0.45), Inches(1.05), Inches(8.0), Inches(5.95))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(245, 247, 250)
    rect.line.color.rgb = RGBColor(130, 140, 150)
    rect.line.width = Pt(1.5)
    tf = rect.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "화면 기획 내용 이미지 삽입 영역\n\n" + (item.screen_summary or "")
    run.font.size = Pt(18)
    run.font.name = "Malgun Gothic"
    run.font.color.rgb = RGBColor(80, 90, 100)


def _add_display_table(slide, item: ScreenPlanItem):
    rows = max(2, len(item.display_items) + 1)
    table_shape = slide.shapes.add_table(rows, 2, Inches(8.75), Inches(1.05), Inches(4.4), Inches(5.95))
    table = table_shape.table
    table.columns[0].width = Inches(1.45)
    table.columns[1].width = Inches(2.95)
    for c, header in enumerate(["표시항목", "설명"]):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(217, 234, 247)
        for paragraph in cell.text_frame.paragraphs:
            if paragraph.runs:
                paragraph.runs[0].font.bold = True
                paragraph.runs[0].font.size = Pt(11)
                paragraph.runs[0].font.name = "Malgun Gothic"
    for r, display_item in enumerate(item.display_items, start=1):
        table.cell(r, 0).text = display_item.item_name
        table.cell(r, 1).text = display_item.description
        for c in range(2):
            for paragraph in table.cell(r, c).text_frame.paragraphs:
                if paragraph.runs:
                    paragraph.runs[0].font.size = Pt(9)
                    paragraph.runs[0].font.name = "Malgun Gothic"


def save_screen_plan_ppt(items: List[ScreenPlanItem], output_path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = title_slide.shapes.add_textbox(Inches(0.6), Inches(2.5), Inches(12), Inches(1))
    _set_text(title_box, "화면기획서", font_size=36, bold=True)
    sub_box = title_slide.shapes.add_textbox(Inches(0.65), Inches(3.45), Inches(12), Inches(0.5))
    _set_text(sub_box, f"총 {len(items)}개 화면", font_size=18)
    if not items:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(0.6), Inches(3.0), Inches(12), Inches(0.8))
        _set_text(box, "화면 관련 요구사항이 추출되지 않았습니다.", font_size=24, bold=True)
    else:
        for item in items:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _add_header(slide, item)
            _add_image_placeholder(slide, item)
            _add_display_table(slide, item)
    prs.save(output_path)
