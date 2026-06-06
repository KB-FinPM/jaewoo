from copy import deepcopy
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from modules.mapper_loader import build_context, build_placeholder_values, get_value, resolve_path
from modules.schemas import ScreenPlanItem
from modules.token_limiter import limit_text_for_output


def _replace_text(text: str, values: Dict[str, str]) -> str:
    for key, value in values.items():
        text = text.replace(key, value or '')
    return text


def _copy_run_style(source_run, target_run):
    """Copy the run properties XML so a newly-created run keeps the template font."""
    source_rpr = source_run._r.get_or_add_rPr()
    target_rpr = target_run._r.get_or_add_rPr()
    target_rpr.clear()
    for child in source_rpr:
        target_rpr.append(deepcopy(child))


def _get_paragraph_end_rpr(paragraph):
    """Return a:endParaRPr from a paragraph. Used when a table cell has no run yet."""
    return paragraph._p.find(qn('a:endParaRPr'))


def _apply_rpr_to_run(run, rpr_source):
    """
    Copy paragraph endParaRPr to a newly-created run as a:rPr.
    Without this, python-pptx may create text with the default Office font.
    """
    if rpr_source is None:
        return

    run_element = run._r
    existing_rpr = run_element.find(qn('a:rPr'))
    if existing_rpr is not None:
        run_element.remove(existing_rpr)

    copied_rpr = deepcopy(rpr_source)
    copied_rpr.tag = qn('a:rPr')
    run_element.insert(0, copied_rpr)


def _set_text_preserve_format(text_frame, new_text: str, force_paragraph_end_style: bool = False, max_token: int = 0):
    """
    Set text without using text_frame.text or cell.text so the template font remains.

    For Description table cells, the template often has no actual run in empty cells.
    In that case, copy a:endParaRPr from the paragraph to the new run's a:rPr
    before writing the text.
    """
    paragraphs = list(text_frame.paragraphs)
    if not paragraphs:
        return

    first_paragraph = paragraphs[0]
    first_runs = list(first_paragraph.runs)

    end_rpr = _get_paragraph_end_rpr(first_paragraph)

    if first_runs:
        if force_paragraph_end_style and end_rpr is not None:
            _apply_rpr_to_run(first_runs[0], end_rpr)
        first_runs[0].text = limit_text_for_output(new_text or '', max_token=max_token)
        for run in first_runs[1:]:
            run.text = ''
    else:
        new_run = first_paragraph.add_run()
        _apply_rpr_to_run(new_run, end_rpr)
        new_run.text = limit_text_for_output(new_text or '', max_token=max_token)

    for paragraph in paragraphs[1:]:
        for run in list(paragraph.runs):
            run.text = ''

def _replace_text_in_text_frame(text_frame, values: Dict[str, str], max_token: int = 0):
    original_text = text_frame.text
    replaced_text = _replace_text(original_text, values)
    if replaced_text != original_text:
        _set_text_preserve_format(text_frame, replaced_text, max_token=max_token)


def _replace_placeholders_in_shape(shape, values: Dict[str, str], max_token: int = 0):
    if getattr(shape, 'has_table', False):
        for row in shape.table.rows:
            for cell in row.cells:
                original_text = cell.text
                replaced_text = _replace_text(original_text, values)
                if replaced_text != original_text:
                    _set_text_preserve_format(cell.text_frame, replaced_text, max_token=max_token)
        return

    if getattr(shape, 'shape_type', None) == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _replace_placeholders_in_shape(child, values, max_token=max_token)
        return

    if getattr(shape, 'has_text_frame', False):
        _replace_text_in_text_frame(shape.text_frame, values, max_token=max_token)


def _replace_placeholders_in_slide(slide, values: Dict[str, str], max_token: int = 0):
    for shape in slide.shapes:
        _replace_placeholders_in_shape(shape, values, max_token=max_token)


def _duplicate_slide(prs, slide_index: int):
    source = prs.slides[slide_index]
    dest = prs.slides.add_slide(source.slide_layout)

    for shape in source.shapes:
        new_element = deepcopy(shape.element)
        dest.shapes._spTree.insert_element_before(new_element, 'p:extLst')

    return dest


def _delete_slide(prs, slide_index: int):
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)
    if slide_index < len(slides):
        slide_id_list.remove(slides[slide_index])


def _iter_tables_from_shapes(shapes):
    for shape in shapes:
        if getattr(shape, 'has_table', False):
            yield shape.table
        if getattr(shape, 'shape_type', None) == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_tables_from_shapes(shape.shapes)


def _find_description_table(slide, table_mapper: Dict[str, Any]):
    header_text = str(table_mapper.get('header_text', '')).strip()
    min_rows = int(table_mapper.get('min_rows', 2))
    min_columns = int(table_mapper.get('min_columns', 2))

    candidates = []
    for table in _iter_tables_from_shapes(slide.shapes):
        if len(table.rows) < min_rows or len(table.columns) < min_columns:
            continue
        first_cell_text = table.cell(0, 0).text.strip()
        if not header_text or first_cell_text == header_text:
            candidates.append(table)

    if candidates:
        return candidates[-1]
    return None


def _format_display_item(display_item: Any, text_format: str) -> str:
    item_name = str(get_value(display_item, 'item_name') or '')
    description = str(get_value(display_item, 'description') or '')
    return text_format.replace('{item_name}', item_name).replace('{description}', description).strip(': ')


def _fill_description_table(slide, item: ScreenPlanItem, table_mapper: Dict[str, Any], max_token: int = 0):
    table = _find_description_table(slide, table_mapper)
    if table is None:
        return

    start_row = int(table_mapper.get('start_row', 1))
    target_column = int(table_mapper.get('target_column', 1))
    max_items = int(table_mapper.get('max_items', 10))
    text_format = table_mapper.get('text_format', '{item_name}: {description}')

    if table_mapper.get('clear_rows_before_fill', True):
        for row_idx in range(start_row, min(len(table.rows), start_row + max_items)):
            _set_text_preserve_format(table.cell(row_idx, target_column).text_frame, '', force_paragraph_end_style=True, max_token=max_token)

    display_items = get_value(item, table_mapper.get('display_items_field', 'display_items')) or []
    for offset, display_item in enumerate(display_items[:max_items]):
        row_idx = start_row + offset
        if row_idx >= len(table.rows):
            break
        description_text = _format_display_item(display_item, text_format)
        _set_text_preserve_format(table.cell(row_idx, target_column).text_frame, description_text, force_paragraph_end_style=True, max_token=max_token)


def save_screen_plan_ppt(
    items: List[ScreenPlanItem],
    template_path: str,
    output_path: str,
    project_name: str,
    author: str,
    mapper: Optional[Dict[str, Any]] = None,
    max_token: int = 0,
):
    mapper = mapper or {
        'template_path': template_path,
        'common_slide_indices': [0, 1],
        'template_slide_index': 2,
        'placeholder_slides': {
            'common': {'{프로젝트명}': 'project_name', '{작성자명}': 'author', '{작성자}': 'author', '{작성일}': 'today'},
            'screen_item': {
                '{요구사항ID}': 'requirement_id',
                '{화면ID}': 'screen_id|screen_no',
                '{화면번호}': 'screen_no',
                '{화면명}': 'screen_name',
                '{서브시스템명}': '',
                '{메뉴위치}': '',
            },
        },
        'description_table': {
            'header_text': 'Description',
            'min_rows': 11,
            'min_columns': 2,
            'start_row': 1,
            'target_column': 1,
            'max_items': 10,
            'display_items_field': 'display_items',
            'text_format': '{item_name}: {description}',
        },
    }

    prs = Presentation(resolve_path(mapper.get('template_path') or template_path))
    context = build_context(project_name, author)

    common_values = build_placeholder_values(
        mapper.get('placeholder_slides', {}).get('common', {}),
        source=None,
        context=context,
    )
    for slide_index in mapper.get('common_slide_indices', [0, 1]):
        if slide_index < len(prs.slides):
            _replace_placeholders_in_slide(prs.slides[slide_index], common_values, max_token=max_token)

    template_slide_index = int(mapper.get('template_slide_index', 2))
    screen_placeholders = mapper.get('placeholder_slides', {}).get('screen_item', {})
    table_mapper = mapper.get('description_table', {})

    for item in items:
        slide = _duplicate_slide(prs, template_slide_index)
        item_values = build_placeholder_values(screen_placeholders, source=item, context=context)
        _replace_placeholders_in_slide(slide, item_values, max_token=max_token)
        _fill_description_table(slide, item, table_mapper, max_token=max_token)

    _delete_slide(prs, template_slide_index)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
