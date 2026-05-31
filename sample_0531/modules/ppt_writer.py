from copy import deepcopy
from datetime import datetime
from pathlib import Path
from typing import Dict, List
import unicodedata

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from modules.schemas import ScreenPlanItem


def _resolve_path(path: str) -> str:
    """Resolve paths even when Korean filenames are stored as NFD on macOS ZIPs."""
    p = Path(path)
    if p.exists():
        return str(p)

    parent = p.parent if str(p.parent) != '' else Path('.')
    target = unicodedata.normalize('NFC', p.name)
    if parent.exists():
        for candidate in parent.iterdir():
            if unicodedata.normalize('NFC', candidate.name) == target:
                return str(candidate)
    raise FileNotFoundError(f'템플릿 파일을 찾을 수 없습니다: {path}')


def _replace_text(text: str, values: Dict[str, str]) -> str:
    for key, value in values.items():
        text = text.replace(key, value or '')
    return text


def _get_paragraph_end_rpr(paragraph):
    """빈 셀에 정의된 endParaRPr을 run 서식으로 복사하기 위해 찾는다."""
    return paragraph._p.find(qn('a:endParaRPr'))


def _apply_rpr_to_run(run, rpr_source):
    """
    새로 만든 run에는 템플릿 폰트가 자동 적용되지 않을 수 있다.
    paragraph의 endParaRPr을 a:rPr로 바꿔 run에 직접 붙여준다.
    """
    if rpr_source is None:
        return

    run_element = run._r

    # 기존 rPr이 있으면 제거
    existing_rpr = run_element.find(qn('a:rPr'))
    if existing_rpr is not None:
        run_element.remove(existing_rpr)

    copied_rpr = deepcopy(rpr_source)
    copied_rpr.tag = qn('a:rPr')
    run_element.insert(0, copied_rpr)


def _set_text_preserve_format(text_frame, new_text: str):
    """
    text_frame.text 또는 cell.text에 직접 대입하면 템플릿 폰트가 초기화될 수 있다.
    기존 run이 있으면 첫 번째 run의 서식을 유지하고 텍스트만 바꾼다.

    Description 표처럼 셀이 비어 있어 run이 없는 경우에는 문단의 endParaRPr
    정보를 새 run의 rPr로 복사해서 템플릿 폰트가 유지되도록 한다.
    """
    paragraphs = list(text_frame.paragraphs)
    if not paragraphs:
        return

    first_paragraph = paragraphs[0]
    first_runs = list(first_paragraph.runs)

    if first_runs:
        first_runs[0].text = new_text or ''
        for run in first_runs[1:]:
            run.text = ''
    else:
        end_rpr = _get_paragraph_end_rpr(first_paragraph)
        new_run = first_paragraph.add_run()
        _apply_rpr_to_run(new_run, end_rpr)
        new_run.text = new_text or ''

    # 첫 번째 문단을 제외한 나머지 문단의 텍스트만 비운다.
    # 문단/런 자체는 유지해서 기존 서식 구조가 최대한 보존되도록 한다.
    for paragraph in paragraphs[1:]:
        for run in list(paragraph.runs):
            run.text = ''


def _replace_text_in_text_frame(text_frame, values: Dict[str, str]):
    # placeholder가 여러 run으로 쪼개져도 전체 문장 기준으로 치환하되,
    # 텍스트 대입은 run.text만 사용해서 템플릿 폰트를 유지한다.
    original_text = text_frame.text
    replaced_text = _replace_text(original_text, values)
    if replaced_text != original_text:
        _set_text_preserve_format(text_frame, replaced_text)


def _replace_placeholders_in_shape(shape, values: Dict[str, str]):
    if getattr(shape, 'has_table', False):
        for row in shape.table.rows:
            for cell in row.cells:
                original_text = cell.text
                replaced_text = _replace_text(original_text, values)
                if replaced_text != original_text:
                    _set_text_preserve_format(cell.text_frame, replaced_text)
        return

    if getattr(shape, 'shape_type', None) == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _replace_placeholders_in_shape(child, values)
        return

    if getattr(shape, 'has_text_frame', False):
        _replace_text_in_text_frame(shape.text_frame, values)


def _replace_placeholders_in_slide(slide, values: Dict[str, str]):
    for shape in slide.shapes:
        _replace_placeholders_in_shape(shape, values)


def _duplicate_slide(prs, slide_index: int):
    source = prs.slides[slide_index]
    # 일부 템플릿은 slide_layouts 개수가 적을 수 있으므로 원본 슬라이드의 레이아웃을 사용한다.
    dest = prs.slides.add_slide(source.slide_layout)

    for shape in source.shapes:
        new_element = deepcopy(shape.element)
        dest.shapes._spTree.insert_element_before(new_element, 'p:extLst')

    return dest


def _delete_slide(prs, slide_index: int):
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)
    if slide_index < len(slides):
        slide_id_list.remove(slides[slide_index])


def _fill_description_table(slide, item: ScreenPlanItem):
    # The screen template has two tables. The Description table is the one with 11 rows and 2 columns.
    description_tables = []
    for shape in slide.shapes:
        if not getattr(shape, 'has_table', False):
            continue
        table = shape.table
        if len(table.columns) >= 2 and len(table.rows) >= 2:
            first_cell_text = table.cell(0, 0).text.strip()
            if first_cell_text == 'Description' or len(table.rows) >= 11:
                description_tables.append(table)

    if not description_tables:
        return

    table = description_tables[-1]
    for row_idx in range(1, min(len(table.rows), 11)):
        _set_text_preserve_format(table.cell(row_idx, 1).text_frame, '')

    for idx, display_item in enumerate(item.display_items[:10], start=1):
        if idx >= len(table.rows):
            break
        description_text = f'{display_item.item_name}: {display_item.description}'.strip(': ')
        _set_text_preserve_format(table.cell(idx, 1).text_frame, description_text)


def save_screen_plan_ppt(
    items: List[ScreenPlanItem],
    template_path: str,
    output_path: str,
    project_name: str,
    author: str,
):
    prs = Presentation(_resolve_path(template_path))
    today = datetime.today().strftime('%Y-%m-%d')

    common_values = {
        '{프로젝트명}': project_name,
        '{작성자명}': author,
        '{작성자}': author,
        '{작성일}': today,
    }

    for slide_index in [0, 1]:
        if slide_index < len(prs.slides):
            _replace_placeholders_in_slide(prs.slides[slide_index], common_values)

    template_slide_index = 2
    for item in items:
        slide = _duplicate_slide(prs, template_slide_index)
        screen_id = item.screen_id or item.screen_no
        item_values = {
            '{요구사항ID}': item.requirement_id,
            '{화면ID}': screen_id,
            '{화면번호}': item.screen_no,
            '{화면명}': item.screen_name,
            '{서브시스템명}': '',
            '{메뉴위치}': '',
        }
        _replace_placeholders_in_slide(slide, item_values)
        _fill_description_table(slide, item)

    _delete_slide(prs, template_slide_index)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
